# personaai/documents/utils.py
import pdfplumber
import openai
import uuid
import re
from docx import Document
from pptx import Presentation
from django.conf import settings
from documents.models import DocumentChunk
from langchain_text_splitters import RecursiveCharacterTextSplitter
from .embedding_utils import get_embedding_model, get_chroma_collection

openai.api_key = settings.OPENAI_API_KEY

def process_document(document_instance):
    collection = get_chroma_collection()
    embedding_model = get_embedding_model()
    doc_path = document_instance.file.path
    text = extract_text(doc_path)
    text = preprocess_text(text)
    chunks = chunk_text(text)
    
    
    embeddings = embedding_model.embed_documents(chunks)
    
    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        embedding_id = f"{document_instance.id}_chunk_{i}_{uuid.uuid4().hex[:8]}"
        collection.add(
            documents=[chunk],
            embeddings=[embedding],
            ids=[embedding_id],
            metadatas=[{"document_id": document_instance.id, "chunk_index": i}]
        )
        DocumentChunk.objects.create(
            document=document_instance,
            chunk_index=i,
            content=chunk,
            embedding_id=embedding_id
        )

def chunk_text(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1024,
        chunk_overlap=200,
        add_start_index=True,
    )
    return text_splitter.split_text(text)

def extract_text_from_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
        return text.strip()

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text).strip()

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read().strip()
    
def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")

def preprocess_text(text):
    text = re.sub(r'[^a-zA-Z0-9\s]', '', text)
    text = text.lower()
    text = re.sub(r'\.{3,}', '.', text)
    text = re.sub(r'\s+', ' ', text)
    while '\n\n' in text:
        text = text.replace('\n\n', '\n')
    return text



# # fungsi delete collection
# def delete_collection(collection):
#     # Dapatkan semua ID dalam collection
#     all_ids = collection.get()["ids"]

#     # Hapus semua dokumen berdasarkan ID
#     collection.delete(ids=all_ids)
#     # hentikan sementara program
#     exit(0)
